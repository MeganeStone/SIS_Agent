from pptx import Presentation
from pptx.enum.text import MSO_AUTO_SIZE
from pptx.util import Pt
import time
import re
import os
import gc
import subprocess
import sys
from openai import OpenAI
from openai import APIError, RateLimitError, APITimeoutError, APIConnectionError
# 新增：导入httpx处理客户端配置
import httpx
# ---------------------- LangChain 依赖 ----------------------
from langchain_core.tools import StructuredTool, ToolException
import tempfile
import shutil
from concurrent.futures import ThreadPoolExecutor, as_completed
import threading
from lxml import etree
import copy

# ---------------------- 全局配置（你的默认目录） ----------------------
DEFAULT_INPUT_DIR = r"D:\seki\AI\copilotTest\input"
DEFAULT_OUTPUT_DIR = r"D:\seki\AI\copilotTest\output"
DEFAULT_TARGET_LANG = "日语"  # 默认翻译目标语言
DEFAULT_DELAY = 1.2  # 每次翻译后的延迟，单位秒（可调整，过快可能触发API限速）
MAX_CONTEXT_LENGTH = 2000  # 上下文最大长度（防止Token超限）

# 全局客户端
_client = None
_translation_cache = {}

def translate_ppt_file(file_name: str, source_dir: str = DEFAULT_INPUT_DIR, output_dir: str = DEFAULT_OUTPUT_DIR, 
                       target_lang: str = DEFAULT_TARGET_LANG, delay: float = DEFAULT_DELAY) -> str:
    """
    翻译PPT文件（并行版）：按幻灯片分组并行翻译，每组内顺序翻译，保持上下文连贯。
    """
    global _translation_cache
    _translation_cache = {}  # 清空缓存
    
    input_path = os.path.abspath(os.path.join(source_dir, file_name))
    output_path = os.path.abspath(os.path.join(output_dir, f"{os.path.splitext(file_name)[0]}_{target_lang}{os.path.splitext(file_name)[1]}"))
    
    # 检查文件
    if not os.path.exists(input_path):
        raise ToolException(f"文件不存在: {input_path}")
    if not os.path.exists(output_dir):
        os.makedirs(output_dir)
    
    from pptx import Presentation
    prs = Presentation(input_path)
    total_slides = len(prs.slides)
    
    # 获取系统可用日语字体（用于后续设置）
    jp_fonts = ["Meiryo UI", "MS Gothic", "MS Mincho", "SimSun", "Arial Unicode MS"]
    jp_font = jp_fonts[0]
    for font in jp_fonts:
        try:
            from pptx.util import Font
            Font(name=font)
            jp_font = font
            break
        except:
            continue
    
    # 收集每个幻灯片的数据：幻灯片对象、需要翻译的形状列表、上下文
    slides_data = []
    for i, slide in enumerate(prs.slides):
        ctx = {
            "file_name": file_name,
            "slide_num": i + 1,
            "total_slides": total_slides,
            "translated_segments": [],
            "term_requirements": "TBOX译为TBOX、TSU译为TSU、CAN总线译为CANバス（日语）/CAN Bus（英语），公司名称不翻译"
        }
        shapes_to_translate = []
        # 遍历形状
        for shape in slide.shapes:
            # 文本框
            if hasattr(shape, "text") and shape.text.strip():
                shapes_to_translate.append(shape)
            # 表格
            if shape.has_table:
                for row in shape.table.rows:
                    for cell in row.cells:
                        if cell.text.strip():
                            shapes_to_translate.append(cell)
        slides_data.append((slide, shapes_to_translate, ctx))
    
    # 定义处理单个幻灯片的函数
    def process_slide(slide, shapes, ctx, target_lang, delay, jp_font):
        # 顺序翻译该幻灯片中的所有形状
        for shape in shapes:
            original_text = shape.text
            translated = translate_text(original_text, target_lang, delay, context=ctx)
            shape.text = translated
            # 设置字体（保证显示）
            if hasattr(shape, "text_frame"):
                tf = shape.text_frame
                tf.word_wrap = True
                from pptx.enum.text import MSO_AUTO_SIZE
                tf.auto_size = MSO_AUTO_SIZE.TEXT_TO_FIT_SHAPE
                for para in tf.paragraphs:
                    for run in para.runs:
                        run.font.name = jp_font
                        if run.font.size is None:
                            from pptx.util import Pt
                            run.font.size = Pt(10)
            # 表格单元格也可能有 text_frame
            elif hasattr(shape, "text_frame"):
                tf = shape.text_frame
                tf.word_wrap = True
                from pptx.enum.text import MSO_AUTO_SIZE
                tf.auto_size = MSO_AUTO_SIZE.TEXT_TO_FIT_SHAPE
                for para in tf.paragraphs:
                    for run in para.runs:
                        run.font.name = jp_font
                        if run.font.size is None:
                            from pptx.util import Pt
                            run.font.size = Pt(10)
        return True
    
    # 并行处理每个幻灯片
    from concurrent.futures import ThreadPoolExecutor, as_completed
    with ThreadPoolExecutor(max_workers=5) as executor:
        futures = []
        for slide, shapes, ctx in slides_data:
            future = executor.submit(process_slide, slide, shapes, ctx, target_lang, delay, jp_font)
            futures.append(future)
        # 等待所有任务完成，若有异常会抛出
        for future in as_completed(futures):
            future.result()
    
    # 保存最终文件
    prs.save(output_path)
    
    # 清理
    _translation_cache.clear()
    gc.collect()
    
    return f"PPT翻译完成！输出路径: {output_path}"

def _get_client(force_recreate=False):
    """获取OpenAI客户端（适配DashScope兼容模式）"""
    global _client
    if _client is None or force_recreate:
        api_key = os.environ.get("OPENAI_API_KEY", "sk-10579025107e412983a48273c2ff7d3f")  # 替换成你的API Key
        base_url = os.environ.get("OPENAI_BASE_URL", "https://dashscope.aliyuncs.com/compatible-mode/v1")
        http_client = httpx.Client(
            timeout=httpx.Timeout(10.0, connect=20.0),  # 显式设置超时
            follow_redirects=True,
        )
        _client = OpenAI(api_key=api_key, base_url=base_url, http_client=http_client)
    return _client

def translate_text(text: str, target_lang: str = DEFAULT_TARGET_LANG, delay: float = DEFAULT_DELAY, 
                   context: dict = None) -> str:
    """
    翻译文本，支持上下文参考和重试机制
    :param text: 要翻译的文本
    :param target_lang: 目标语言（默认日语）
    :param delay: 翻译延迟（默认1.2秒）
    :param context: 上下文信息
    :return: 翻译结果
    """
    # 初始化默认上下文
    context = context or {
        "file_name": "未知文件",
        "slide_num": 0,
        "total_slides": 0,
        "translated_segments": [],
        "term_requirements": "汽车行业术语（如TBOX/TSU/CAN总线）需统一翻译，公司名称不翻译"
    }
    if not text or not isinstance(text, str) or text.strip() == "":
        return text
    # 基础过滤
    text = re.sub(r"请提供需要翻译的原文内容.*", "", text)
    if target_lang == "日语" and re.search(r"[\u3040-\u30ff]", text):
        return text
    company_names = (
        "上海畅星", "上海暢星",
        "上海畅星软件有限公司", "上海暢星軟件有限公司",
        "上海畅星软件有限会社", "上海暢星ソフトウェア有限公司"
    )
    for name in company_names:
        if name in text:
            return text
    # if not re.search(r"[\u4e00-\u9fff]", text):
    #     return text
    
    # 缓存key加入上下文标识（保证同文件同页的缓存隔离）
    cache_key = f"{context['file_name']}_{context['slide_num']}_{text}_{target_lang}"
    if cache_key in _translation_cache:
        print(f"[缓存命中] 页码{context['slide_num']} | 文本长度: {len(text)}")
        return _translation_cache[cache_key]

    print(f"\n[开始翻译] 文本长度: {len(text)} | 目标语言: {target_lang}")

    # ====================== 核心修改：构建带上下文的Prompt ======================
    # 拼接已翻译的上下文（只保留最近的，控制长度）
    # translated_context = "\n".join(context["translated_segments"][-3:])  # 只保留最近3段 todo：token充裕时可增加上下文数量，提升连贯性，但要注意Token限制
    translated_context = '暂无已翻译段落'  # 当前token限制较紧，暂时关闭上下文传递，后续可根据实际情况调整是否开启上下文参考
    # 截断上下文，防止Token超限
    if len(translated_context) > MAX_CONTEXT_LENGTH:
        translated_context = translated_context[-MAX_CONTEXT_LENGTH:]

    # 带上下文的系统提示
    system_prompt = f"""
    你是汽车TSU(本公司开发的产品，Telematic Systems Unit)专业翻译官，需遵守以下规则：
    1. 翻译目标：将文本内容翻译成{target_lang}，数字保持不变；
    2. 上下文参考：本次翻译属于文件「{context['file_name']}」的第{context['slide_num']}页（共{context['total_slides']}页）；
    3. 连贯要求：需参考当前页已翻译的段落保持术语、语气、格式统一：
       {translated_context if translated_context else '暂无已翻译段落'}
    4. 术语要求：{context['term_requirements']}；
    5. 输出要求：只输出翻译结果，无任何解释、说明或额外文字。
    """
    
    # 用户提示：明确待翻译文本
    user_prompt = f"请翻译以下内容：\n{text}"
    
    translation = None
    # 重试策略
    for attempt in range(2):
        try:
            client = _get_client(force_recreate=(attempt > 0))
            # 修正：使用正确的OpenAI调用方法（chat.completions.create）
            resp = client.chat.completions.create(
                model="qwen-plus",
                messages=[
                    {"role": "system", "content": system_prompt},
                    {"role": "user", "content": user_prompt}
                ],
                temperature=0.1
            )
            translation = resp.choices[0].message.content.strip()
            # 过滤额外输出
            translation = re.sub(r"^[\s\u3000]*翻译结果：|^[\s\u3000]*译文：", "", translation)
            if translation and (target_lang != "日语" or re.search(r"[\u3040-\u30ff\u31f0-\u31ff]", translation)):
                print(f"[翻译成功] 尝试{attempt+1}次，结果长度: {len(translation)}")
                # 更新上下文：将本次翻译结果加入已翻译段落（供后续段落参考）
                context["translated_segments"].append(f"原文：{text} | 译文：{translation}")
                break
            else:
                print(f"[翻译无效] 尝试{attempt+1}次，结果无{target_lang}内容，重试...")
                translation = None
        except (APIError, RateLimitError, APITimeoutError, APIConnectionError) as e:
            print(f"[重试] 尝试{attempt+1}次失败（API错误）: {e}")
            time.sleep(delay * (attempt + 1))
        except Exception as e:
            print(f"[重试] 尝试{attempt+1}次失败（其他错误）: {e}")
            time.sleep(delay * (attempt + 1))

    # 兜底
    if not translation or translation.strip() == "":
        print(f"[翻译兜底] 所有重试失败，使用原文")
        translation = text
    else:
        time.sleep(delay)

    # 缓存翻译结果
    _translation_cache[cache_key] = translation
    return translation


# ---------------------- 新增：强制清理Excel进程 ----------------------
def kill_excel_processes():
    """杀死所有残留的Excel进程，解决目录锁定问题"""
    try:
        if sys.platform == "win32":
            # Windows杀死Excel进程（兼容中英文系统）
            subprocess.run(["taskkill", "/f", "/im", "EXCEL.EXE"], capture_output=True, encoding="gbk", errors="ignore")
            subprocess.run(["taskkill", "/f", "/im", "excel.exe"], capture_output=True, encoding="gbk", errors="ignore")
        print("已清理所有残留Excel进程")
    except Exception as e:
        print(f"清理Excel进程失败: {e}")

def translate_excel_all_text(wb, sheet_count, target_lang, delay, excel_context):
    """
    完整翻译Excel工作表的所有文字元素（单元格+形状+批注+图表+页眉页脚）
    """
    # 遍历所有工作表（带序号）
    for i, ws in enumerate(wb.Worksheets, 1):
        excel_context["slide_num"] = i
        excel_context["translated_segments"] = []  # 每个Sheet独立上下文
        print(f"\n=== Excel 第 {i}/{sheet_count} 个工作表: {ws.Name} ===")

        # ========== 1. 翻译单元格（优化UsedRange，遍历所有有内容的单元格） ==========
        try:
            # 刷新UsedRange，避免漏选
            ws.UsedRange
            # 遍历所有单元格（替代仅UsedRange，可根据需求调整为UsedRange提升效率）
            # 方式1：高效版（仅UsedRange）
            for cell in ws.UsedRange:
                val = cell.Value
                if val and isinstance(val, str):
                    val_stripped = val.strip()
                    # if re.search(r"[\u4e00-\u9fff]", val_stripped):
                    cell.Value = translate_text(val_stripped, target_lang, delay, context=excel_context)
        
        except Exception as e:
            print(f"翻译单元格失败: {e}")

        # ========== 2. 翻译所有形状文字（含分组/TextFrame2） ==========
        def translate_shape_text(shp):
            """递归翻译形状（含分组形状）的文字"""
            try:
                # 处理分组形状：递归遍历子形状
                if shp.Type == 6:  # 6=xlGroup：分组形状
                    for sub_shp in shp.GroupItems:
                        translate_shape_text(sub_shp)
                    return
                
                # 处理SmartArt图形
                if shp.Type == 19:  # 19=xlSmartArt：SmartArt图形
                    for node in shp.SmartArt.AllNodes:
                        if node.TextFrame2.TextRange.Text:
                            t = node.TextFrame2.TextRange.Text.strip()
                            # if re.search(r"[\u4e00-\u9fff]", t):
                            node.TextFrame2.TextRange.Text = translate_text(t, target_lang, delay, context=excel_context)
                    return

                # 优先用TextFrame2（Office 2007+），兼容TextFrame
                if hasattr(shp, "TextFrame2") and shp.TextFrame2.HasText:
                    t = shp.TextFrame2.TextRange.Text.strip()
                    # if re.search(r"[\u4e00-\u9fff]", t):
                    shp.TextFrame2.TextRange.Text = translate_text(t, target_lang, delay, context=excel_context)
                elif hasattr(shp, "TextFrame") and shp.TextFrame.Characters.Text:
                    t = shp.TextFrame.Characters.Text.strip()
                    # if re.search(r"[\u4e00-\u9fff]", t):
                    shp.TextFrame.Characters().Text = translate_text(t, target_lang, delay, context=excel_context)
            except Exception as e:
                print(f"翻译形状[{shp.Name}]失败: {e}")

        # 遍历所有形状（含分组）
        for shp in ws.Shapes:
            translate_shape_text(shp)

        # ========== 3. 翻译单元格批注/备注 ==========
        try:
            # Excel 2016+用ws.Comments，旧版用ws.Notes
            for comment in ws.Comments:
                t = comment.Text().strip()
                if re.search(r"[\u4e00-\u9fff]", t):
                    comment.Text(translate_text(t, target_lang, delay, context=excel_context))
        except Exception as e:
            print(f"翻译批注失败: {e}")

        # ========== 4. 翻译图表中的文字 ==========
        try:
            for chart in ws.ChartObjects:
                # 翻译图表标题
                if chart.Chart.HasTitle:
                    t = chart.Chart.ChartTitle.Text.strip()
                    # if re.search(r"[\u4e00-\u9fff]", t):
                    chart.Chart.ChartTitle.Text = translate_text(t, target_lang, delay, context=excel_context)
                # 翻译坐标轴标签（X/Y轴）
                for axis in chart.Chart.Axes:
                    if axis.HasTitle:
                        t = axis.AxisTitle.Text.strip()
                        # if re.search(r"[\u4e00-\u9fff]", t):
                        axis.AxisTitle.Text = translate_text(t, target_lang, delay, context=excel_context)
                # 翻译数据标签（可选）
                # for series in chart.Chart.SeriesCollection():
                #     if series.HasDataLabels:
                #         for label in series.DataLabels:
                #             t = label.Text.strip()
                #             if re.search(r"[\u4e00-\u9fff]", t):
                #                 label.Text = translate_text(t, target_lang, delay, context=excel_context)
        except Exception as e:
            print(f"翻译图表文字失败: {e}")

        # ========== 5. 翻译页眉/页脚文字 ==========
        try:
            # 翻译页眉（左/中/右）
            for align in ["LeftHeader", "CenterHeader", "RightHeader"]:
                t = getattr(ws.PageSetup, align).strip()
                # if re.search(r"[\u4e00-\u9fff]", t):
                setattr(ws.PageSetup, align, translate_text(t, target_lang, delay, context=excel_context))
            # 翻译页脚（左/中/右）
            for align in ["LeftFooter", "CenterFooter", "RightFooter"]:
                t = getattr(ws.PageSetup, align).strip()
                # if re.search(r"[\u4e00-\u9fff]", t):
                setattr(ws.PageSetup, align, translate_text(t, target_lang, delay, context=excel_context))
        except Exception as e:
            print(f"翻译页眉页脚失败: {e}")

def translate_excel_file(file_name: str, source_dir: str = DEFAULT_INPUT_DIR, output_dir: str = DEFAULT_OUTPUT_DIR,
                         target_lang: str = DEFAULT_TARGET_LANG, delay: float = DEFAULT_DELAY) -> str:
    """
    翻译Excel文件（保留所有元素：单元格、文本框、图形）
    :param file_name: 要翻译的Excel文件名（如"test.xlsx"）
    :param source_dir: 源文件目录（默认D:\seki\AI\copilotTest\input）
    :param output_dir: 输出目录（默认D:\seki\AI\copilotTest\output）
    :param target_lang: 目标语言（默认日语）
    :param delay: 翻译延迟（默认2.0秒）
    :return: 翻译结果提示
    """
    global _translation_cache
    _translation_cache = {}  # 清空缓存
    input_path = os.path.abspath(os.path.join(source_dir, file_name)).replace("/", "\\")
    output_path = os.path.abspath(os.path.join(output_dir, f"{os.path.splitext(file_name)[0]}_{target_lang}{os.path.splitext(file_name)[1]}")).replace("/", "\\")

    # 检查文件/目录
    if not os.path.exists(input_path):
        raise ToolException(f"Excel文件不存在: {input_path}")
    if not os.path.exists(output_dir):
        os.makedirs(output_dir)
    if not os.access(output_dir, os.W_OK):
        raise ToolException(f"输出目录无写入权限: {output_dir}")

    # 清理Excel进程
    kill_excel_processes()

    try:
        import pythoncom
        import win32com.client as win32
        pythoncom.CoInitialize()

        # 使用DispatchEx避免共享实例
        excel = win32.DispatchEx("Excel.Application")
        excel.Visible = False
        excel.DisplayAlerts = False
        excel.AskToUpdateLinks = False

        wb = excel.Workbooks.Open(
            Filename=input_path,
            ReadOnly=False,
            IgnoreReadOnlyRecommended=True,
            UpdateLinks=0
        )
        sheet_count = wb.Worksheets.Count
        # ====================== 核心修改：初始化Excel级上下文 ======================
        excel_context = {
            "file_name": file_name,
            "slide_num": 0,
            "total_slides": sheet_count,
            "translated_segments": [],  # 存储当前页已翻译的段落
            "term_requirements": "TBOX译为TBOX、TSU译为TSU、CAN总线译为CANバス（日语）/CAN Bus（英语），公司名称不翻译"
        }
        # 调用完整翻译函数
        translate_excel_all_text(wb, sheet_count, target_lang, delay, excel_context)

        # 保存文件（兼容不同Excel版本）
        wb.SaveAs(
            Filename=output_path,
            FileFormat=51,  # 51 = xlOpenXMLWorkbook (xlsx)
            ReadOnlyRecommended=False,
            CreateBackup=False
        )
        # 安全关闭
        wb.Close(SaveChanges=False)
        excel.Quit()
        pythoncom.CoUninitialize()
        # 再次清理进程
        kill_excel_processes()
        # 清理缓存和内存
        _translation_cache.clear()
        gc.collect()
        return f"Excel翻译完成！输出路径: {output_path}"
    except Exception as e:
        # 异常时确保清理进程
        try:
            pythoncom.CoUninitialize()
        except:
            pass
        kill_excel_processes()
        _translation_cache.clear()
        gc.collect()
        raise ToolException(f"Excel翻译失败: {str(e)}")
    
if __name__ == "__main__":
    """独立测试 PPT 翻译工具的入口"""
    import sys
    import traceback
    from pathlib import Path

    # ============== 测试配置 ==============
    TEST_FILE = "test.pptx"  # 测试文件名
    TEST_INPUT_DIR = Path(DEFAULT_INPUT_DIR)
    TEST_OUTPUT_DIR = Path(DEFAULT_OUTPUT_DIR)
    
    # ============== 验证文件存在 ==============
    input_path = TEST_INPUT_DIR / TEST_FILE
    if not input_path.exists():
        print(f"❌ 错误：测试文件不存在！\n路径: {input_path}\n请确保文件已上传到 {TEST_INPUT_DIR}")
        print("\n请检查：")
        print("1. 文件名是否包含特殊字符（如 ⊿）")
        print("2. 文件是否在正确目录：", TEST_INPUT_DIR.resolve())
        print("3. 文件扩展名是否为 .pptx")
        sys.exit(1)
    
    print(f"✅ 文件存在: {input_path.resolve()}")
    print(f"✅ 测试目录: {TEST_INPUT_DIR.resolve()}")
    print(f"✅ 输出目录: {TEST_OUTPUT_DIR.resolve()}\n")
    
    # ============== 执行翻译测试 ==============
    print("="*50)
    print("🚀 开始翻译测试...")
    print("="*50)
    
    try:
        # 清空缓存（确保干净测试）
        _translation_cache.clear()
        
        # 执行翻译
        start_time = time.time()
        result = translate_ppt_file(
            file_name=TEST_FILE,
            source_dir=DEFAULT_INPUT_DIR,
            output_dir=DEFAULT_OUTPUT_DIR,
            target_lang="英语",
            delay=DEFAULT_DELAY
        )
        elapsed = time.time() - start_time
        
        # ============== 结果验证 ==============
        print("\n" + "="*50)
        print("✅ 翻译完成！")
        print(f"耗时: {elapsed:.2f} 秒")
        print(f"输出文件: {Path(result).resolve()}")
        print("="*50)
        
        # 检查输出文件
        output_path = Path(DEFAULT_OUTPUT_DIR) / f"{TEST_FILE.split('.')[0]}_{DEFAULT_TARGET_LANG}.pptx"
        if output_path.exists():
            print(f"✓ 输出文件已生成: {output_path.resolve()}")
            print(f"✓ 文件大小: {output_path.stat().st_size} 字节")
        else:
            print(f"❌ 输出文件未生成: {output_path.resolve()}")
            print("提示：可能因路径特殊字符导致文件未创建")
        
    except Exception as e:
        print(f"\n❌ 翻译测试失败: {str(e)}")
        print("详细错误:")
        traceback.print_exc()
        sys.exit(1)
    
    print("\n" + "="*50)
    print("测试完成！")
    print("="*50)